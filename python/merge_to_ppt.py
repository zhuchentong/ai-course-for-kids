import os
from pptx import Presentation
from pptx.util import Inches

SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

blank_layout = prs.slide_layouts[6]

images_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), "images")
files = sorted(f for f in os.listdir(images_dir) if f.endswith(".jpeg"))

for filename in files:
    slide = prs.slides.add_slide(blank_layout)
    pic = slide.shapes.add_picture(
        os.path.join(images_dir, filename),
        0, 0, SLIDE_WIDTH, SLIDE_HEIGHT
    )

output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "output.pptx")
prs.save(output_path)
print(f"Done: {len(files)} slides -> {output_path}")
